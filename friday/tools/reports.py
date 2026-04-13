"""
reports.py — כלי יצירת מצגות וסיכומים
יוצר קבצי PPTX ודוחות טקסטואליים בעברית
"""
from pptx import Presentation
from pptx.util import Inches, Pt
import os
from friday.tools._client import get_anthropic_client

def register(mcp):
    @mcp.tool()
    async def generate_presentation(topic: str, points: str, output_path: str = "presentation.pptx") -> str:
        """יוצר מצגת PPTX בעברית. topic=נושא, points=נקודות עיקריות מופרדות בפסיק"""
        client = get_anthropic_client()

        # בקש מ-Claude לפרק לשקפים
        response = client.messages.create(
            model="claude-sonnet-4-6",
            max_tokens=1000,
            messages=[{
                "role": "user",
                "content": f"""צור מבנה מצגת עסקית בעברית על הנושא: {topic}
נקודות עיקריות: {points}

החזר JSON בפורמט הבא בלבד (ללא טקסט נוסף):
{{
  "title": "כותרת המצגת",
  "slides": [
    {{"title": "כותרת שקף", "bullets": ["נקודה 1", "נקודה 2", "נקודה 3"]}},
    ...
  ]
}}

צור 5-6 שקפים מקצועיים."""
            }]
        )

        import json
        raw = response.content[0].text.strip()
        raw = raw.replace("```json", "").replace("```", "").strip()
        data = json.loads(raw)

        # בנה את ה-PPTX
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # שקף כותרת
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = data["title"]
        title_slide.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(36)
        title_slide.shapes.title.text_frame.paragraphs[0].runs[0].font.bold = True

        # שקפי תוכן
        for slide_data in data["slides"]:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = slide_data["title"]
            tf = slide.placeholders[1].text_frame
            tf.clear()
            for i, bullet in enumerate(slide_data["bullets"]):
                p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                p.text = bullet
                p.level = 0

        prs.save(output_path)
        import subprocess
        subprocess.Popen(
            ['powershell', '-c', f'Start-Process "{os.path.abspath(output_path)}"'],
            shell=True
        )
        return f"Presentation created and opened: {os.path.abspath(output_path)} ({len(data['slides'])} slides)"

    @mcp.tool()
    async def generate_html_presentation(topic: str, points: str) -> str:
        """יוצר מצגת HTML מלאה עם dark theme שנפתחת בדפדפן. topic=נושא, points=נקודות מופרדות בפסיק"""
        import json, tempfile, subprocess, datetime

        client = get_anthropic_client()

        # בקש מ-Claude מבנה שקפים
        response = client.messages.create(
            model="claude-sonnet-4-6",
            max_tokens=1500,
            messages=[{
                "role": "user",
                "content": f"""צור מבנה מצגת עסקית מקצועית בעברית על הנושא: {topic}
נקודות עיקריות: {points}

החזר JSON בלבד (ללא טקסט נוסף):
{{
  "title": "כותרת ראשית",
  "subtitle": "תת-כותרת או שם המציג",
  "slides": [
    {{"title": "כותרת שקף", "bullets": ["נקודה 1", "נקודה 2", "נקודה 3"]}},
    ...
  ],
  "cta": "קריאה לפעולה לשקף האחרון"
}}

צור 4-6 שקפים אמצעיים מקצועיים. כל שקף 3-5 נקודות."""
            }]
        )

        raw = response.content[0].text.strip().replace("```json", "").replace("```", "").strip()
        data = json.loads(raw)

        title = data.get("title", topic)
        subtitle = data.get("subtitle", "")
        slides = data.get("slides", [])
        cta = data.get("cta", "נפגש לדון בשלב הבא")
        total = len(slides) + 2  # title + content slides + summary

        # בנה slide divs
        slide_divs = []

        # שקף ראשון — כותרת
        slide_divs.append(f"""
    <div class="slide active" id="slide-0">
      <div class="slide-content title-slide">
        <div class="title-badge">FRIDAY PRESENTATION</div>
        <h1 class="main-title">{title}</h1>
        <p class="main-subtitle">{subtitle}</p>
        <div class="slide-meta">{datetime.date.today().strftime("%d.%m.%Y")} &nbsp;|&nbsp; {total} שקפים</div>
      </div>
    </div>""")

        # שקפים אמצעיים
        for i, slide in enumerate(slides, start=1):
            bullets_html = "\n".join(
                f'          <li class="bullet-item" style="animation-delay:{j*0.1:.1f}s">{b}</li>'
                for j, b in enumerate(slide.get("bullets", []))
            )
            slide_divs.append(f"""
    <div class="slide" id="slide-{i}">
      <div class="slide-content content-slide">
        <div class="slide-number">{i:02d} / {total:02d}</div>
        <h2 class="slide-title">{slide['title']}</h2>
        <ul class="bullets">
{bullets_html}
        </ul>
      </div>
    </div>""")

        # שקף אחרון — סיכום
        summary_points = [s["title"] for s in slides[:4]]
        summary_html = "\n".join(
            f'          <li class="summary-item">{p}</li>' for p in summary_points
        )
        slide_divs.append(f"""
    <div class="slide" id="slide-{total-1}">
      <div class="slide-content summary-slide">
        <h2 class="summary-heading">סיכום</h2>
        <ul class="summary-list">
{summary_html}
        </ul>
        <div class="cta-box">{cta}</div>
      </div>
    </div>""")

        all_slides = "\n".join(slide_divs)

        html = f"""<!DOCTYPE html>
<html lang="he" dir="rtl">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>{title}</title>
<style>
  * {{ margin: 0; padding: 0; box-sizing: border-box; }}

  body {{
    background: #0a0a0f;
    color: #e0e0e0;
    font-family: 'Segoe UI', Arial, sans-serif;
    height: 100vh;
    overflow: hidden;
    direction: rtl;
  }}

  /* ── slides ── */
  .slide {{
    position: absolute;
    inset: 0;
    display: flex;
    align-items: center;
    justify-content: center;
    opacity: 0;
    transform: translateX(60px);
    transition: opacity .45s ease, transform .45s ease;
    pointer-events: none;
  }}
  .slide.active {{
    opacity: 1;
    transform: translateX(0);
    pointer-events: all;
  }}
  .slide.exit {{
    opacity: 0;
    transform: translateX(-60px);
    transition: opacity .3s ease, transform .3s ease;
  }}

  .slide-content {{
    width: 88%;
    max-width: 1100px;
    padding: 60px 70px;
    background: linear-gradient(135deg, #12121a 0%, #1a1a2e 100%);
    border: 1px solid #2a2a3e;
    border-radius: 20px;
    box-shadow: 0 0 60px rgba(0,200,255,.07);
  }}

  /* ── title slide ── */
  .title-badge {{
    font-size: 11px;
    letter-spacing: 3px;
    color: #00c8ff;
    text-transform: uppercase;
    margin-bottom: 28px;
  }}
  .main-title {{
    font-size: clamp(2rem, 4vw, 3.2rem);
    font-weight: 700;
    color: #ffffff;
    line-height: 1.2;
    margin-bottom: 18px;
    text-shadow: 0 0 30px rgba(0,200,255,.3);
  }}
  .main-subtitle {{
    font-size: 1.2rem;
    color: #8888aa;
    margin-bottom: 36px;
  }}
  .slide-meta {{
    font-size: 0.85rem;
    color: #444466;
    border-top: 1px solid #2a2a3e;
    padding-top: 20px;
  }}

  /* ── content slide ── */
  .slide-number {{
    font-size: 11px;
    letter-spacing: 2px;
    color: #00c8ff55;
    margin-bottom: 20px;
  }}
  .slide-title {{
    font-size: clamp(1.5rem, 2.5vw, 2.2rem);
    font-weight: 700;
    color: #ffffff;
    margin-bottom: 36px;
    padding-bottom: 16px;
    border-bottom: 2px solid #00c8ff33;
  }}
  .bullets {{
    list-style: none;
    display: flex;
    flex-direction: column;
    gap: 18px;
  }}
  .bullet-item {{
    display: flex;
    align-items: flex-start;
    gap: 14px;
    font-size: clamp(1rem, 1.4vw, 1.2rem);
    color: #ccccdd;
    opacity: 0;
    transform: translateY(10px);
    animation: fadeUp .4s forwards;
  }}
  .bullet-item::before {{
    content: '›';
    color: #00c8ff;
    font-size: 1.4em;
    line-height: 1;
    flex-shrink: 0;
    margin-top: 2px;
  }}
  @keyframes fadeUp {{
    to {{ opacity: 1; transform: translateY(0); }}
  }}

  /* ── summary slide ── */
  .summary-heading {{
    font-size: 2rem;
    color: #ffffff;
    margin-bottom: 30px;
    padding-bottom: 14px;
    border-bottom: 2px solid #00c8ff33;
  }}
  .summary-list {{
    list-style: none;
    display: grid;
    grid-template-columns: 1fr 1fr;
    gap: 16px;
    margin-bottom: 36px;
  }}
  .summary-item {{
    background: #1e1e30;
    border: 1px solid #2a2a45;
    border-radius: 10px;
    padding: 16px 20px;
    font-size: 1rem;
    color: #aaaacc;
    border-right: 3px solid #00c8ff;
  }}
  .cta-box {{
    background: linear-gradient(90deg, #003344, #001a22);
    border: 1px solid #00c8ff44;
    border-radius: 12px;
    padding: 22px 28px;
    font-size: 1.1rem;
    color: #00c8ff;
    text-align: center;
    font-weight: 600;
  }}

  /* ── navigation ── */
  #nav {{
    position: fixed;
    bottom: 32px;
    left: 50%;
    transform: translateX(-50%);
    display: flex;
    align-items: center;
    gap: 16px;
    background: #111118ee;
    border: 1px solid #2a2a3e;
    border-radius: 40px;
    padding: 10px 22px;
    backdrop-filter: blur(10px);
    z-index: 100;
  }}
  .nav-btn {{
    background: none;
    border: 1px solid #2a2a3e;
    color: #aaaacc;
    width: 38px;
    height: 38px;
    border-radius: 50%;
    cursor: pointer;
    font-size: 1.1rem;
    transition: all .2s;
    display: flex; align-items: center; justify-content: center;
  }}
  .nav-btn:hover {{ background: #00c8ff22; border-color: #00c8ff66; color: #fff; }}
  #slide-counter {{ font-size: 13px; color: #666688; min-width: 55px; text-align: center; }}

  /* ── progress bar ── */
  #progress {{
    position: fixed;
    top: 0; left: 0;
    height: 3px;
    background: linear-gradient(90deg, #00c8ff, #7b2fff);
    transition: width .4s ease;
    z-index: 100;
  }}

  /* ── fullscreen btn ── */
  #fs-btn {{
    position: fixed;
    top: 18px;
    left: 18px;
    background: #111118cc;
    border: 1px solid #2a2a3e;
    color: #666688;
    border-radius: 8px;
    padding: 7px 12px;
    cursor: pointer;
    font-size: 13px;
    transition: all .2s;
    z-index: 100;
  }}
  #fs-btn:hover {{ color: #fff; border-color: #00c8ff55; }}
</style>
</head>
<body>

<div id="progress"></div>
<button id="fs-btn" onclick="toggleFS()">&#x26F6; מסך מלא</button>

{all_slides}

<div id="nav">
  <button class="nav-btn" onclick="go(-1)">&#8249;</button>
  <span id="slide-counter">1 / {total}</span>
  <button class="nav-btn" onclick="go(1)">&#8250;</button>
</div>

<script>
  const total = {total};
  let cur = 0;

  function go(dir) {{
    const slides = document.querySelectorAll('.slide');
    const next = cur + dir;
    if (next < 0 || next >= total) return;
    slides[cur].classList.remove('active');
    slides[cur].classList.add('exit');
    setTimeout(() => slides[cur].classList.remove('exit'), 350);
    cur = next;
    slides[cur].classList.add('active');
    document.getElementById('slide-counter').textContent = (cur + 1) + ' / ' + total;
    document.getElementById('progress').style.width = ((cur + 1) / total * 100) + '%';
  }}

  document.addEventListener('keydown', e => {{
    if (e.key === 'ArrowRight' || e.key === 'ArrowDown') go(-1);
    if (e.key === 'ArrowLeft'  || e.key === 'ArrowUp')   go(1);
  }});

  function toggleFS() {{
    if (!document.fullscreenElement) document.documentElement.requestFullscreen();
    else document.exitFullscreen();
  }}

  // init progress
  document.getElementById('progress').style.width = (1 / total * 100) + '%';
</script>
</body>
</html>"""

        # שמור ופתח
        out_path = os.path.abspath(f"presentation_{topic[:20].replace(' ', '_')}.html")
        with open(out_path, "w", encoding="utf-8") as f:
            f.write(html)

        subprocess.Popen(
            ['powershell', '-c', f'Start-Process "{out_path}"'],
            shell=True
        )
        return f"HTML presentation created and opened: {out_path} ({total} slides)"

    @mcp.tool()
    async def generate_summary_doc(text: str) -> str:
        """מסכם טקסט ארוך לדוח עסקי מסודר בעברית"""
        client = get_anthropic_client()
        response = client.messages.create(
            model="claude-sonnet-4-6",
            max_tokens=600,
            messages=[{
                "role": "user",
                "content": f"""סכם את הטקסט הבא לדוח עסקי מסודר בעברית:

{text}

פורמט הסיכום:
- נושא מרכזי
- נקודות מפתח (3-5)
- מסקנות
- המלצות לפעולה"""
            }]
        )
        return response.content[0].text
